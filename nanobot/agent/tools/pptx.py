"""PowerPoint operations tools using python-pptx."""

import json
import re
from pathlib import Path
from typing import Any

from nanobot.agent.tools.base import Tool


def _clean_text(text: str) -> str:
    """Clean text by removing control characters and fixing encoding issues."""
    if not text:
        return ""
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)
    return text.strip()


def _resolve_path(path: str, workspace: Path | None = None, allowed_dir: Path | None = None) -> Path:
    """Resolve path against workspace and enforce directory restriction."""
    p = Path(path).expanduser()
    if not p.is_absolute() and workspace:
        p = workspace / p
    resolved = p.resolve()
    if allowed_dir:
        try:
            resolved.relative_to(allowed_dir.resolve())
        except ValueError:
            raise PermissionError(f"Path {path} is outside allowed directory {allowed_dir}")
    return resolved


class PptxReadTool(Tool):
    """Tool to read PowerPoint presentation."""

    def __init__(self, workspace: Path | None = None, allowed_dir: Path | None = None):
        self._workspace = workspace
        self._allowed_dir = allowed_dir

    @property
    def name(self) -> str:
        return "pptx_read"

    @property
    def description(self) -> str:
        return "Read PowerPoint presentation structure and text content."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "The absolute path to the PowerPoint file (.pptx)"
                },
                "extract_text": {
                    "type": "boolean",
                    "description": "Extract text from slides. Default: true"
                }
            },
            "required": ["path"]
        }

    async def execute(self, path: str, extract_text: bool = True, **kwargs: Any) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return "Error: python-pptx is required. Install with: pip install python-pptx"

        try:
            file_path = _resolve_path(path, self._workspace, self._allowed_dir)
            if not file_path.exists():
                return f"Error: File not found: {path}"
            if file_path.suffix.lower() not in [".pptx", ".pptm"]:
                return f"Error: Not a PowerPoint file: {path}"

            prs = Presentation(str(file_path))

            slides_data = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_info = {"slide_number": slide_num}

                if extract_text:
                    text_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text = _clean_text(shape.text)
                            if text:
                                text_content.append(text)
                    slide_info["text"] = "\n".join(text_content)

                slides_data.append(slide_info)

            result = {
                "slide_count": len(slides_data),
                "slides": slides_data
            }

            return json.dumps(result, ensure_ascii=False, indent=2)

        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            import traceback
            return f"Error reading PowerPoint: {str(e)}\n{traceback.format_exc()}"


class PptxCreateTool(Tool):
    """Tool to create PowerPoint presentation."""

    def __init__(self, workspace: Path | None = None, allowed_dir: Path | None = None):
        self._workspace = workspace
        self._allowed_dir = allowed_dir

    @property
    def name(self) -> str:
        return "pptx_create"

    @property
    def description(self) -> str:
        return "Create a new PowerPoint presentation with title and content."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "output_path": {
                    "type": "string",
                    "description": "The absolute path for the output PowerPoint file"
                },
                "title": {
                    "type": "string",
                    "description": "Presentation title"
                },
                "slides": {
                    "type": "array",
                    "description": "Array of slide objects with 'title' and 'content' keys"
                }
            },
            "required": ["output_path", "title"]
        }

    async def execute(self, output_path: str, title: str, slides: list | None = None, **kwargs: Any) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return "Error: python-pptx is required. Install with: pip install python-pptx"

        try:
            out_path = _resolve_path(output_path, self._workspace, self._allowed_dir)

            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            title_shape.text = title
            subtitle_shape.text = "Generated by AI"

            if slides:
                bullet_slide_layout = prs.slide_layouts[1]
                for slide_data in slides:
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]

                    title_shape.text = slide_data.get("title", "")

                    tf = body_shape.text_frame
                    tf.text = ""

                    content = slide_data.get("content", "")
                    if isinstance(content, str):
                        for line in content.split("\n"):
                            p = tf.add_paragraph()
                            p.text = line
                    elif isinstance(content, list):
                        for item in content:
                            p = tf.add_paragraph()
                            p.text = item

            out_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(out_path))

            return f"Successfully created presentation with {len(slides) + 1 if slides else 1} slides: {out_path}"

        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error creating PowerPoint: {str(e)}"


class PptxAddSlideTool(Tool):
    """Tool to add a slide to PowerPoint presentation."""

    def __init__(self, workspace: Path | None = None, allowed_dir: Path | None = None):
        self._workspace = workspace
        self._allowed_dir = allowed_dir

    @property
    def name(self) -> str:
        return "pptx_add_slide"

    @property
    def description(self) -> str:
        return "Add a new slide to an existing PowerPoint presentation."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "The absolute path to the PowerPoint file"
                },
                "title": {
                    "type": "string",
                    "description": "Slide title"
                },
                "content": {
                    "type": "string",
                    "description": "Slide content (bullets, one per line)"
                },
                "layout": {
                    "type": "string",
                    "enum": ["title", "bullet", "blank", "title_only"],
                    "description": "Slide layout type. Default: bullet"
                }
            },
            "required": ["path", "title"]
        }

    async def execute(self, path: str, title: str, content: str | None = None,
                    layout: str = "bullet", **kwargs: Any) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return "Error: python-pptx is required. Install with: pip install python-pptx"

        try:
            file_path = _resolve_path(path, self._workspace, self._allowed_dir)
            if not file_path.exists():
                return f"Error: File not found: {path}"

            prs = Presentation(str(file_path))

            layout_map = {
                "title": 0,
                "bullet": 1,
                "blank": 6,
                "title_only": 5
            }
            layout_index = layout_map.get(layout, 1)

            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)

            if layout in ["title", "bullet", "title_only"]:
                if slide.shapes.title:
                    slide.shapes.title.text = title

            if content and layout == "bullet":
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = ""
                for line in content.split("\n"):
                    if line.strip():
                        p = tf.add_paragraph()
                        p.text = line.strip()

            prs.save(str(file_path))

            return f"Successfully added slide to {file_path}"

        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error adding slide: {str(e)}"


class PptxExtractImagesTool(Tool):
    """Tool to extract images from PowerPoint presentation."""

    def __init__(self, workspace: Path | None = None, allowed_dir: Path | None = None):
        self._workspace = workspace
        self._allowed_dir = allowed_dir

    @property
    def name(self) -> str:
        return "pptx_extract_images"

    @property
    def description(self) -> str:
        return "Extract all images from a PowerPoint presentation."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "The absolute path to the PowerPoint file"
                },
                "output_dir": {
                    "type": "string",
                    "description": "Output directory to save extracted images"
                }
            },
            "required": ["path", "output_dir"]
        }

    async def execute(self, path: str, output_dir: str, **kwargs: Any) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return "Error: python-pptx is required. Install with: pip install python-pptx"

        try:
            file_path = _resolve_path(path, self._workspace, self._allowed_dir)
            if not file_path.exists():
                return f"Error: File not found: {path}"

            out_path = _resolve_path(output_dir, self._workspace, self._allowed_dir)
            out_path.mkdir(parents=True, exist_ok=True)

            prs = Presentation(str(file_path))
            image_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        ext = image.ext
                        image_name = f"slide{slide_num}_{image_count + 1}.{ext}"
                        image_path = out_path / image_name

                        with open(image_path, "wb") as f:
                            f.write(image_bytes)
                        image_count += 1

            return f"Successfully extracted {image_count} images to {out_path}"

        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error extracting images: {str(e)}"


class PptxListSlidesTool(Tool):
    """Tool to list slides in PowerPoint presentation."""

    def __init__(self, workspace: Path | None = None, allowed_dir: Path | None = None):
        self._workspace = workspace
        self._allowed_dir = allowed_dir

    @property
    def name(self) -> str:
        return "pptx_list_slides"

    @property
    def description(self) -> str:
        return "List all slides in a PowerPoint presentation with their titles."

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "The absolute path to the PowerPoint file"
                }
            },
            "required": ["path"]
        }

    async def execute(self, path: str, **kwargs: Any) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return "Error: python-pptx is required. Install with: pip install python-pptx"

        try:
            file_path = _resolve_path(path, self._workspace, self._allowed_dir)
            if not file_path.exists():
                return f"Error: File not found: {path}"

            prs = Presentation(str(file_path))

            slides_info = []
            for slide_num, slide in enumerate(prs.slides, 1):
                title = ""
                if slide.shapes.title and slide.shapes.title.text:
                    title = _clean_text(slide.shapes.title.text)

                slides_info.append({
                    "slide_number": slide_num,
                    "title": title
                })

            result = {
                "slide_count": len(slides_info),
                "slides": slides_info
            }

            return json.dumps(result, ensure_ascii=False, indent=2)

        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error listing slides: {str(e)}"
